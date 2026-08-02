import os
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import pytesseract
import google.generativeai as genai
from app.config import settings

# Automatically detect Windows common install paths for Tesseract OCR if not in PATH
POSSIBLE_TESSERACT_PATHS = [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    os.path.expanduser(r"~\AppData\Local\Programs\Tesseract-OCR\tesseract.exe")
]

for t_path in POSSIBLE_TESSERACT_PATHS:
    if os.path.exists(t_path):
        pytesseract.pytesseract.tesseract_cmd = t_path
        break

def extract_text_from_image(file_path: str) -> str:
    """Extracts text from image files using Tesseract OCR with Gemini Vision API fallback."""
    extracted_text = ""

    # Attempt 1: Tesseract OCR
    try:
        image = Image.open(file_path)
        ocr_result = pytesseract.image_to_string(image)
        if ocr_result and len(ocr_result.strip()) > 5:
            return ocr_result.strip()
    except Exception as tesseract_err:
        print(f"[Tesseract OCR Warning] {tesseract_err}")

    # Attempt 2: Gemini Vision API Fallback
    if settings.GEMINI_API_KEY and len(settings.GEMINI_API_KEY) > 10:
        try:
            genai.configure(api_key=settings.GEMINI_API_KEY)
            # Use gemini-1.5-flash or gemini-pro-vision for OCR image processing
            model = genai.GenerativeModel("gemini-1.5-flash")
            image = Image.open(file_path)
            response = model.generate_content([
                "Extract and transcribe all readable text, titles, headings, and numbers from this image verbatim. Do not summarize or format as code.",
                image
            ])
            if response and response.text:
                return response.text.strip()
        except Exception as gemini_err:
            print(f"[Gemini Vision OCR Error] {gemini_err}")

    return extracted_text.strip() or "No readable text could be detected in this image file."

def extract_text_from_file(file_path: str, file_type: str) -> str:
    """Extracts raw text from PDF, DOCX, PPTX, TXT, or Image files."""
    if not os.path.exists(file_path):
        return ""

    file_type = file_type.lower()
    text = ""

    try:
        if file_type == "pdf":
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text() + "\n"
        elif file_type in ["docx", "doc"]:
            doc = DocxDocument(file_path)
            for p in doc.paragraphs:
                text += p.text + "\n"
        elif file_type in ["pptx", "ppt"]:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_type in ["txt", "text"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif file_type in ["png", "jpg", "jpeg", "webp", "bmp"]:
            text = extract_text_from_image(file_path)
    except Exception as e:
        print(f"[File Extraction Error] Failed to extract text: {e}")
        text = "Could not extract full text from file."

    return text.strip() or "No text could be extracted from this document."
